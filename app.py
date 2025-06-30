import requests
import re
import speech_recognition as sr
import os 
import streamlit as st
from pptx import Presentation
from playsound import playsound
from gtts import gTTS
from io import BytesIO
from sentence_transformers import SentenceTransformer
from sklearn.metrics.pairwise import cosine_similarity
import numpy as np

r = sr.Recognizer()

# API OpenRouter
url = "https://openrouter.ai/api/v1/chat/completions"
headers = {
    "Authorization": f"Bearer {os.getenv('DEEPSEEK_API_KEY')}",
    "Content-Type": "application/json",
}

def extract_text_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    text_list = []
    
    for slide in prs.slides:
        text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
        text_list.append(text)
    
    return "\n".join(text_list)  # Combined all text into one long string

def chunking(study_material):
    max_len = 1000 # Max text length in each chunk
    sentences = re.split(r'(?<=[.!?]) +', study_material)
    chunks = []
    curr_chunk = ""

    for s in sentences:
        if len(curr_chunk) + len(s) <= max_len:
            curr_chunk += " " + s
        else:
            chunks.append(curr_chunk.strip())
            curr_chunk = s # Start a new chunk with a new text
    
    if curr_chunk:
        chunks.append(curr_chunk.strip())

    return chunks

def find_chunk(chunks, question):
    # Find the most relevant chunk using Sentence embedding
    model = SentenceTransformer("all-MiniLM-L6-v2") 
    embeddings = model.encode(chunks + [question])
    
    chunk_embeddings = embeddings[:-1]
    question_embedding = embeddings[-1].reshape(1, -1)

    similarities = cosine_similarity(question_embedding, chunk_embeddings)[0]
    best_chunk_index = np.argmax(similarities)

    return chunks[best_chunk_index]

def ask_deepseek(question, context, topic):
    all_folder = "data/"
    study_material = ""

    # Pick the selected topic
    topic_folders = {
        "Machine Learning": "ML",
        "Artificial Intelligence": "AI",
    }

    # Read the slides only from the selected topic
    if topic == "All Topics":
        selected_folders = [os.path.join(all_folder, folder) for folder in topic_folders.values()]
    else:
        selected_folders = [os.path.join(all_folder, topic_folders[topic])]

    # Extract text from ppt yg from selected folder
    for folder in selected_folders:
        if os.path.exists(folder):  
            for file in os.listdir(folder):
                if file.endswith(".pptx"):  
                    ppts = os.path.join(folder, file)
                    study_material += extract_text_from_ppt(ppts) + "\n\n"

    # Chunking the text and find the most relevant text (from the ppt) based on the user's question  
    chunks = chunking(study_material)
    most_relevant_chunk = find_chunk(chunks, question)

    # Additional validation
    if not isinstance(most_relevant_chunk, str) or not most_relevant_chunk.strip():
        most_relevant_chunk = "No relevant information found."

    # Structured Prompt that pass the most relevant chunk and the user's question to the DeepSeek AI to generate final response
    data = {
        "model": "deepseek/deepseek-r1:free",
        "messages": [
            {
                "role": "system",
                "content": (
                    "If you are writing any math formulas, note that you have a MathJax render environment.\n"
                    "- Any LaTeX text between single dollar sign ($) will be rendered as a TeX formula.\n"
                    "- Use $(tex_formula)$ in-line delimiters to display equations instead of backslash.\n"
                    "- The render environment only uses $ (single dollar sign) as a container delimiter, never output $$.\n"
                    "Example: $x^2 + 3x$ is output for 'x² + 3x' to appear as TeX.\n\n"
                    "Answer the questions based only on the provided context."
                ),
            },

            {"role": "system", 
             "content": (
                "Answer the user's questions as a friendly AI assistant.\n"
                "- If the question is related to the given context, provide an informative and relevant answer.\n"
                "- If the question is general or not related to the context, respond naturally like a conversational AI. Be friendly.\n"  
                "- Be concise, but engaging and friendly in your responses.\n"
                "- If unsure, let the user know you don't have enough context but try to help anyway.\n"  
             )
            },
            {"role": "user", 
             "content": f"Here is some study material:\n\n{most_relevant_chunk}\n\nBased on the above content, answer concisely: {question}"}
        ]
    }

    # Debugging
    print("Thinking...")

    response = requests.post(url, headers=headers, json=data)

    # Check response
    if response.status_code == 200:
        return response.json().get("choices", [{}])[0].get("message", {}).get("content", "No response.")
    
    return "Error fetching response."

def textToSpeech(response):
    if not response or response.strip() == "":  # Check response whether it's empty or not
        st.error("Deepseek mengembalikan respons kosong.")
        response = "Maaf, saya tidak bisa menemukan jawaban."
        return response
    
    audio_bytes = BytesIO()
    tts = gTTS(text=response, lang="en") # Make the speech audio from the response text 
    tts.write_to_fp(audio_bytes)
    audio_bytes.seek(0)

    return audio_bytes.read()

def speechToText():
    try:
        with sr.Microphone() as source2: # Activating microphone 
            print("Mic On! I'm all ears! :)")
            st.toast("🎤 Mic On! You can speak now")
            print("")

            r.adjust_for_ambient_noise(source2, duration=0.2)
            audio2 = r.listen(source2) 

            # Convert speech to text in English
            prompt = r.recognize_google(audio2, language="en-US")
            prompt = prompt.lower()

            print("Your question here: ", prompt)

            st.session_state.messages.append({"role": "user", "content": prompt})

            st.toast("✅ Voice detected!")
            
            with st.chat_message("user"):
                st.write(prompt)

            # Get AI response
            response = ask_deepseek(prompt, st.session_state.context, st.session_state.topic)
            cleaned_response = re.sub(r"[*_#]", "", response)  # Clean markdown format
            audio_data = textToSpeech(cleaned_response)

            # Save response to chat history
            st.session_state.messages.append({"role": "assistant", "content": cleaned_response, "audio": audio_data if audio_data else None})
            st.session_state.audio_files.append(audio_data)

            with st.chat_message("assistant"):
                st.write(cleaned_response)
                st.audio(audio_data, format="audio/mp3")

    except Exception as e:
        print(e)
        print("Could not request results; {0}".format(e))

# GUI via streamlit
st.set_page_config(page_title="Chatbot", layout="wide")

with st.sidebar:
    col1, col2 = st.columns([0.4, 0.85]) 

    with col1:
        with st.container(): 
            st.image("logo-cleaned.png", width=75)  

    with col2:
        st.markdown("""
            <div style="display: flex; flex-direction: column; justify-content: center; height: 100%;">
                <h1 style="margin: 0; padding: 0;">BISA</h1>
                <h4 style="margin: 0; margin-bottom: 10px;">BINUS Intelligent Student Assistant</h3>
            </div>
        """, unsafe_allow_html=True)

with st.sidebar:
    st.markdown(
        """
        <p style='text-align: left; font-size: 14px; color: gray; margin-bottom: 30px;'>
        BISA is an AI-powered student assistant designed to help BINUS students with academic questions, especially in Computer Science field. You can select a topic for more specific answers, or just ask anything directly!
        </p>
        """, 
        unsafe_allow_html=True
    )

    topics = st.sidebar.selectbox(
        "Select Topic", 
        ["All Topics", "Machine Learning", "Artificial Intelligence"], 
        index=0
    )

    st.session_state.topic = topics

    st.markdown(
        """
        <p style='text-align: left; font-size: 14px; color: black; margin-bottom: 5px;'>
            Try asking using your voice!
        </p>
        """, 
        unsafe_allow_html=True
    )

    if st.button("🎤 Speech Mode", key="mic_button", help="Click to start voice input"):
        st.session_state.speech_mode = True 

if "messages" not in st.session_state:
    st.session_state.messages = [] 
    st.session_state.audio_files = []
    st.session_state.context = ""

    # AI Greets user when the chat starts for the first time
    opening_greet = "Hi! I'm BISA, your student assistant. Anything I can help you with?"
    opening_audio = textToSpeech(opening_greet)

    st.session_state.messages = [{"role": "assistant", "content": opening_greet, "audio": opening_audio}]
    st.session_state.audio_files.append(opening_audio)

for msg in st.session_state.messages:
    with st.chat_message(msg["role"]):
        st.write(msg["content"])
        if msg["role"] == "assistant" and "audio" in msg:
            st.audio(msg["audio"], format="audio/mp3")

if st.session_state.get("speech_mode", False):
    speechToText()  
    st.session_state.speech_mode = False
    st.rerun()

if user_input := st.chat_input("Ask your question..."):
    st.session_state.messages.append({"role": "user", "content": user_input})

    with st.chat_message("user"):
        st.write(user_input)

    response = ask_deepseek(user_input, st.session_state.context, topics)
    cleaned_response = re.sub(r"[*_#]", "", response)   
    audio_data = textToSpeech(cleaned_response)

    st.session_state.messages.append({"role": "assistant", "content": response, "audio": audio_data if audio_data else None})
    st.session_state.audio_files.append(audio_data)

    with st.chat_message("assistant"):
        st.write(response)
        st.audio(audio_data, format="audio/mp3")